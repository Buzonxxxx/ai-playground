from dotenv import load_dotenv
from openai import OpenAI
import pandas as pd
import os
import time

load_dotenv()
openai_api_key = os.getenv("OPENAI_API_KEY")

client = OpenAI()

# Read file
file_path = 'sales_data.csv'
sales_data = pd.read_csv(file_path)

# Create doc
file = client.files.create(
    file=open(file_path, 'rb'),  # read binary
    purpose='assistants'
)
# Create an agent with doc
assistant = client.beta.assistants.create(
    name="數據分析師",
    instructions="作為一位數據科學助理，當給定數據和一個查詢時，你能編寫適當的程式碼並創建適當的視覺化。",
    model="gpt-4o",
    tools=[{ "type": "code_interpreter" }],
    tool_resources={ "code_interpreter": { "file_ids": [file.id] }}
)
print(assistant)
print("----------")

# Create thread
thread = client.beta.threads.create(
    messages=[
        {
            "role": "user",
            "content": "計算從2022年到2025年每個季度的總銷售總額，並通過不同的產品將其可視化為折線圖，產品線條顏色分別為紅，藍，綠。",
            "attachments": [
                {
                    "file_id": file.id,
                    "tools": [{ "type": "code_interpreter" }]
                }
            ]
        }
    ]
)
print(thread)
print("----------")

# Create run
run = client.beta.threads.runs.create(
    thread_id=thread.id,
    assistant_id=assistant.id,
)
print(run)
print("----------")

# Wait run to finish
while True:
    messages = client.beta.threads.messages.list(thread_id=thread.id)
    try:
        messages.data[0].content[0].image_file
        time.sleep(5)
        print("圖表已建立")
        if messages.data and messages.data[0].content:
            print("當前Message: ", messages.data[0].content[0])
        break
    except:
        time.sleep(10)
        print("你的助手正在努力的做圖表....")
        if messages.data and messages.data[0].content:
            print("當前Message: ", messages.data[0].content[0])

# Convert output file to png and save it
def convert_file_to_png(file_id, write_path):
    data = client.files.content(file_id)
    data_bytes = data.read()
    with open(write_path, "wb") as file:
        file.write(data_bytes)

plot_file_id = messages.data[0].content[0].image_file.file_id
image_path = "圖書銷售.png"
convert_file_to_png(plot_file_id, image_path)

# Upload png file
plot_file = client.files.create(
  file=open(image_path, "rb"),
  purpose='assistants'
)

# messages = client.beta.threads.messages.list(thread_id=thread.id)
# [message.content[0] for message in messages.data]

# Define submit message function
def submit_message_wait_completion(assistant_id, thread, user_message, file_ids=None):
    # 檢查並等待所有正在執行的Run完成
    for run in client.beta.threads.runs.list(thread_id=thread.id).data:
        if run.status == 'in_progress':
            print(f"等待Run {run.id} 完成...")
            while True:
                run_status = client.beta.threads.runs.retrieve(thread_id=thread.id, run_id=run.id).status
                if run_status in ['succeeded', 'failed']:
                    break
                time.sleep(5)  # wait for 5s then check again

    params = {
        'thread_id': thread.id,
        'role': 'user',
        'content': user_message,
    }
    if file_ids:        
        attachments = [{"file_id": file_id, "tools": [ {"type": "code_interpreter"}]} for file_id in file_ids]
        params['attachments'] = attachments
    client.beta.threads.messages.create(**params)

    run = client.beta.threads.runs.create(thread_id=thread.id, assistant_id=assistant_id)
    return run 


def get_response(thread):
    return client.beta.threads.messages.list(thread_id=thread.id)

submit_message_wait_completion(assistant.id, thread, "請根據你剛才建立的圖表，给我兩個约20字的句子，描述最重要的洞察。這會用來做投影片展示，揭露出數據背後的'秘密'。")
time.sleep(10)
response = get_response(thread)
bullet_points = response.data[0].content[0].text.value
print(bullet_points)
print("----------")

submit_message_wait_completion(assistant.id, thread, "根據你創建的情節和要點，為投影片想一個非常簡短的標題。它應該只反映你得出的主要見解。")
time.sleep(10)
response = get_response(thread)
title = response.data[0].content[0].text.value
print(title)
print("----------")

# 提供花语秘境公司的说明
company_summary = "我们是網路鲜花批發商，但是我们董事长也寫IT圖書！"

# 使用DALL-E 3生成圖片
response = client.images.generate(
  model='dall-e-3',
  prompt=f"根据這個公司概述 {company_summary}, \
           生成一張展示成長和前進道路的啟發性照片。這將用於季度銷售規劃會議",
       size="1024x1024",
       quality="hd",
       n=1
)
image_url = response.data[0].url

# 獲取DALL-E 3生成的圖片
import requests
dalle_img_path = '花語秘境.png'
img = requests.get(image_url)

# 將圖片存到本地
with open(dalle_img_path,'wb') as file:
  file.write(img.content)

# 上傳圖片提供给助手做為PPT素材
dalle_file = client.files.create(
  file=open(dalle_img_path, "rb"),
  purpose='assistants'
)

title_template = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# 创建新的演示文稿对象
prs = Presentation()

# 添加一个空白的幻灯片布局
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# 将幻灯片的背景颜色设置为黑色
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# 在幻灯片左侧添加图片，上下留有边距
left = Inches(0)
top = Inches(0)
height = prs.slide_height
width = prs.slide_width * 3/5
pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)

# 在较高位置添加标题文本框
left = prs.slide_width * 3/5
top = Inches(2)
width = prs.slide_width * 2/5
height = Inches(1)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.bold = True
title_p.font.size = Pt(38)
title_p.font.color.rgb = RGBColor(255, 255, 255)
title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# 添加副标题文本框
left = prs.slide_width * 3/5
top = Inches(3)
width = prs.slide_width * 2/5
height = Inches(1)
subtitle_box = slide.shapes.add_textbox(left, top, width, height)
subtitle_frame = subtitle_box.text_frame
subtitle_p = subtitle_frame.add_paragraph()
subtitle_p.text = subtitle_text
subtitle_p.font.size = Pt(22)
subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
subtitle_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
"""

data_vis_template = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# 创建新的演示文稿对象
prs = Presentation()

# 添加一个空白的幻灯片布局
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# 将幻灯片的背景颜色设置为黑色
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# 定义占位符
image_path = data_vis_img
title_text = "提升利润：在线销售与直销优化的主导地位"
bullet_points = "• 在线销售在各个季度中始终领先于盈利能力，表明了强大的数字市场存在。\n• 直销表现出波动，表明该渠道的表现变化和需要针对性改进的必要性。"

# 在幻灯片左侧添加图片占位符
left = Inches(0.2)
top = Inches(1.8)
height = prs.slide_height - Inches(3)
width = prs.slide_width * 3/5
pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)

# 添加覆盖整个宽度的标题文本
left = Inches(0)
top = Inches(0)
width = prs.slide_width
height = Inches(1)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.margin_top = Inches(0.1)
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.bold = True
title_p.font.size = Pt(28)
title_p.font.color.rgb = RGBColor(255, 255, 255)
title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# 添加硬编码的“关键见解”文本和项目符号列表
left = prs.slide_width * 2/3
top = Inches(1.5)
width = prs.slide_width * 1/3
height = Inches(4.5)
insights_box = slide.shapes.add_textbox(left, top, width, height)
insights_frame = insights_box.text_frame
insights_p = insights_frame.add_paragraph()
insights_p.text = "关键见解："
insights_p.font.bold = True
insights_p.font.size = Pt(24)
insights_p.font.color.rgb = RGBColor(0, 128, 100)
insights_p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
insights_frame.add_paragraph()


bullet_p = insights_frame.add_paragraph()
bullet_p.text = bullet_points
bullet_p.font.size = Pt(12)
bullet_p.font.color.rgb = RGBColor(255, 255, 255)
bullet_p.line_spacing = 1.5
"""

title_text = "花语秘境"
subtitle_text = "2025年銷售大會"

submit_message_wait_completion(assistant.id,thread,f"使用包含的代碼模板創建符合模板格式的PPTX投影片，但使用本消息中包含的圖片、公司名稱/標题和文件名/副標题：\
{title_template}。重要提示：在此第一张投影片中使用本消息中包含的圖片文件作为image_path圖像，並使用公司名稱 {title_text} 作爲title_text變量，\
  使用副標题文本 {subtitle_text} 作为subtitle_text變量。\
    接着，使用以下代碼模板创建第二张投影片：{data_vis_template}，創建符合模板格式的PPTX投影片，但使用公司名稱/標题和文件名/副標题：\
{data_vis_template}。重要提示：使用您之前在本thread中創建的第二張附图（折線图）作为data_vis_img圖像，並使用您之前創建的數據可視化標题作爲title_text變量，\
  使用您之前創建的見解项目符號列表作为bullet_points變量。将這两张投影片輸出为.pptx文件。確保輸出爲兩張投影片，每張投影片都符合本消息中给出的相應模板。",
              file_ids=[dalle_file.id, plot_file.id]
)


# 等待助手完成PPT创建任务
while True:
    try:
        response = get_response(thread)
        pptx_id = response.data[0].content[0].text.annotations[0].file_path.file_id
        print("成功檢索到 pptx_id:", pptx_id)
        break
    except Exception as e:
        print("您的Assistant正在努力製作投影片...")
        time.sleep(10)

import io
pptx_id = response.data[0].content[0].text.annotations[0].file_path.file_id
ppt_file= client.files.content(pptx_id)
file_obj = io.BytesIO(ppt_file.read())
with open("花语秘境.pptx", "wb") as f:
    f.write(file_obj.getbuffer())